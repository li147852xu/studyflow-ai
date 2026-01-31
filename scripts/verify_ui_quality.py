from __future__ import annotations

import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from core.ingest.cite import build_citation
from infra.db import get_workspaces_dir
from infra.models import init_db
from service.document_service import count_documents, list_documents, set_document_type
from service.ingest_service import ingest_document
from service.tasks_service import enqueue_ingest_index_task, run_task_by_id
from service.workspace_service import create_workspace


def _run_cmd(args: list[str]) -> None:
    result = subprocess.run(args, check=False)
    if result.returncode != 0:
        raise SystemExit(result.returncode)


def _verify_imports() -> None:
    __import__("app.main")


def _build_sample_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Sample DOCX", level=1)
    doc.add_paragraph("This is a short paragraph for ingestion.")
    doc.save(path)


def _build_sample_pptx(path: Path) -> None:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX"
    slide.placeholders[1].text = "This slide contains sample text."
    presentation.save(path)


def _verify_library_ingest() -> None:
    workspace_id = create_workspace("verify-ui-quality")
    uploads_dir = get_workspaces_dir() / workspace_id / "uploads"
    uploads_dir.mkdir(parents=True, exist_ok=True)

    pdf_path = Path(__file__).resolve().parents[1] / "examples" / "ml_fundamentals.pdf"
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_dir_path = Path(temp_dir)
        docx_path = temp_dir_path / "sample.docx"
        pptx_path = temp_dir_path / "sample.pptx"
        _build_sample_docx(docx_path)
        _build_sample_pptx(pptx_path)

        for path, doc_type in [
            (pdf_path, "paper"),
            (docx_path, "course"),
            (pptx_path, "other"),
        ]:
            ingest_document(
                workspace_id=workspace_id,
                filename=path.name,
                data=path.read_bytes(),
                save_dir=uploads_dir,
                ocr_mode="auto",
                ocr_threshold=50,
                doc_type=doc_type,
            )

    docs = list_documents(
        workspace_id,
        sort_by="created_at",
        sort_order="desc",
        limit=10,
        offset=0,
    )
    if not docs:
        raise SystemExit("No documents ingested.")
    for doc in docs:
        if not doc.get("created_at"):
            raise SystemExit("Document missing created_at.")

    docs_by_size = list_documents(
        workspace_id,
        sort_by="size_bytes",
        sort_order="asc",
        limit=10,
        offset=0,
    )
    if docs_by_size and docs_by_size[0].get("size_bytes") is None:
        raise SystemExit("size_bytes missing on documents.")

    first_doc = docs[0]
    set_document_type(doc_id=first_doc["id"], doc_type="other")
    if count_documents(workspace_id, doc_type="other") <= 0:
        raise SystemExit("doc_type update failed.")


def _verify_task_flow() -> None:
    workspace_id = create_workspace("verify-ui-quality-tasks")
    uploads_dir = get_workspaces_dir() / workspace_id / "uploads"
    uploads_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = Path(__file__).resolve().parents[1] / "examples" / "ml_fundamentals.pdf"
    task_id = enqueue_ingest_index_task(
        workspace_id=workspace_id,
        path=str(pdf_path),
        ocr_mode="off",
        ocr_threshold=50,
        doc_type="paper",
        save_dir=str(uploads_dir),
        write_file=True,
    )
    result = run_task_by_id(task_id)
    if not result.get("ingest"):
        raise SystemExit("Ingest task did not return result.")


def _verify_citations() -> None:
    citation = build_citation(
        filename="sample.docx",
        page_start=2,
        page_end=3,
        text="Sample citation text.",
        file_type="docx",
    )
    if "para" not in citation.location_label:
        raise SystemExit("Citation location label missing for DOCX.")


def main() -> None:
    init_db()
    _run_cmd([sys.executable, "-m", "compileall", "."])
    if shutil.which("ruff"):
        _run_cmd(["ruff", "check", "."])
    _run_cmd([sys.executable, "-m", "pytest", "-q"])
    _verify_imports()
    _verify_library_ingest()
    _verify_task_flow()
    _verify_citations()
    print("verify_ui_quality: OK")


if __name__ == "__main__":
    main()
