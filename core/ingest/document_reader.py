from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation

from core.ingest.ocr import OCRSettings, ocr_available, run_ocr
from core.ingest.pdf_reader import PDFPage


class DocumentReadError(RuntimeError):
    pass


@dataclass
class DocumentParseResult:
    pages: list[PDFPage]
    warnings: list[str] = field(default_factory=list)

    @property
    def page_count(self) -> int:
        return len(self.pages)


def _read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="utf-8", errors="ignore")


def read_text_lines(path: Path) -> DocumentParseResult:
    if not path.exists():
        raise DocumentReadError("Text file not found.")
    text = _read_text(path)
    lines = text.splitlines()
    pages = [
        PDFPage(
            number=index,
            text=line,
            text_source="extract",
            ocr_text=None,
            image_count=0,
            has_images=False,
            blocks_json=None,
        )
        for index, line in enumerate(lines, start=1)
        if line.strip()
    ]
    if not pages:
        raise DocumentReadError("Text file has no readable content.")
    return DocumentParseResult(pages=pages)


def read_docx(path: Path) -> DocumentParseResult:
    if not path.exists():
        raise DocumentReadError("DOCX file not found.")
    try:
        document = DocxDocument(str(path))
    except Exception as exc:
        raise DocumentReadError("Failed to open DOCX file.") from exc
    pages: list[PDFPage] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        pages.append(
            PDFPage(
                number=index,
                text=text,
                text_source="extract",
                ocr_text=None,
                image_count=0,
                has_images=False,
                blocks_json=None,
            )
        )
    if not pages:
        raise DocumentReadError("DOCX contains no readable paragraphs.")
    return DocumentParseResult(pages=pages)


def read_pptx(path: Path) -> DocumentParseResult:
    if not path.exists():
        raise DocumentReadError("PPTX file not found.")
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DocumentReadError("Failed to open PPTX file.") from exc
    pages: list[PDFPage] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                parts.append(text)
        slide_text = "\n".join(parts).strip()
        if not slide_text:
            continue
        pages.append(
            PDFPage(
                number=slide_index,
                text=slide_text,
                text_source="extract",
                ocr_text=None,
                image_count=0,
                has_images=False,
                blocks_json=None,
            )
        )
    if not pages:
        raise DocumentReadError("PPTX contains no readable text.")
    return DocumentParseResult(pages=pages)


def read_html(path: Path) -> DocumentParseResult:
    if not path.exists():
        raise DocumentReadError("HTML file not found.")
    html = _read_text(path)
    soup = BeautifulSoup(html, "html.parser")
    elements = soup.find_all(
        ["p", "li", "h1", "h2", "h3", "h4", "h5", "h6", "pre", "blockquote"]
    )
    pages: list[PDFPage] = []
    for index, element in enumerate(elements, start=1):
        text = " ".join(element.stripped_strings).strip()
        if not text:
            continue
        pages.append(
            PDFPage(
                number=index,
                text=text,
                text_source="extract",
                ocr_text=None,
                image_count=0,
                has_images=False,
                blocks_json=None,
            )
        )
    if not pages:
        raise DocumentReadError("HTML contains no readable text elements.")
    return DocumentParseResult(pages=pages)


def read_image(
    path: Path, *, ocr_mode: str = "auto", ocr_settings: OCRSettings | None = None
) -> DocumentParseResult:
    if not path.exists():
        raise DocumentReadError("Image file not found.")
    ocr_settings = ocr_settings or OCRSettings()
    ocr_ready, ocr_reason = ocr_available(ocr_settings)
    if ocr_mode == "off":
        raise DocumentReadError("OCR is disabled for images.")
    if not ocr_ready:
        raise DocumentReadError(f"OCR unavailable: {ocr_reason}")
    try:
        image = Image.open(path)
    except Exception as exc:
        raise DocumentReadError("Failed to open image file.") from exc
    ocr_text = run_ocr(image, settings=ocr_settings).strip()
    if not ocr_text:
        raise DocumentReadError("No text detected in the image.")
    pages = [
        PDFPage(
            number=1,
            text=ocr_text,
            text_source="ocr",
            ocr_text=ocr_text,
            image_count=1,
            has_images=True,
            blocks_json=None,
        )
    ]
    return DocumentParseResult(pages=pages)
